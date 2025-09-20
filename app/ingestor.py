import os
import fitz  # PyMuPDF
import docx
from bs4 import BeautifulSoup
from pptx import Presentation
from sqlalchemy import create_engine, Column, Integer, String, Text
from sqlalchemy.orm import declarative_base
from sqlalchemy.orm import sessionmaker
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from sqlalchemy.exc import OperationalError


def get_user_by_username(username):
    return session.query(User).filter(User.username == username).first()


Base = declarative_base()


def create_user(username, email, full_name, hashed_password, is_admin=0):
    user = User(
        username=username,
        email=email,
        full_name=full_name,
        hashed_password=hashed_password,
        disabled=0,
        is_admin=is_admin
    )
    session.add(user)
    session.commit()
    return user


DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:////data/app.db")
engine = create_engine(DATABASE_URL)
try:
    # Testa conexão inicial; se falhar, usa SQLite local
    with engine.connect() as _:
        pass
except Exception:
    fallback_url = "sqlite:////data/app.db"
    engine = create_engine(fallback_url)
Session = sessionmaker(bind=engine)
session = Session()

Base = declarative_base()


# Modelo de usuário para autenticação
class User(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True)
    username = Column(String, unique=True, nullable=False)
    email = Column(String, unique=True, nullable=False)
    full_name = Column(String)
    hashed_password = Column(String, nullable=False)
    disabled = Column(Integer, default=0)
    is_admin = Column(Integer, default=0)


class ConteudoTeologico(Base):
    __tablename__ = "conteudo_teologico"
    id = Column(Integer, primary_key=True)
    titulo = Column(String)
    texto = Column(Text)
    tipo = Column(String)
    autor = Column(String)
    tema = Column(String)
    fonte = Column(String)


Base.metadata.create_all(engine)


def extrair_pdf(caminho):
    doc = fitz.open(caminho)
    texto = "\n".join([pagina.get_text() for pagina in doc])
    return texto


def extrair_docx(caminho):
    doc = docx.Document(caminho)
    return "\n".join([p.text for p in doc.paragraphs])


def extrair_html(caminho):
    with open(caminho, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
        return soup.get_text()


def extrair_txt(caminho):
    with open(caminho, "r", encoding="utf-8") as f:
        return f.read()


def extrair_ppt(caminho):
    prs = Presentation(caminho)
    texto = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto.append(shape.text)
    return "\n".join(texto)


def salvar_conteudo(titulo, texto, tipo, autor=None, tema=None, fonte=None):
    conteudo = ConteudoTeologico(
        titulo=titulo,
        texto=texto,
        tipo=tipo,
        autor=autor,
        tema=tema,
        fonte=fonte
    )
    session.add(conteudo)
    session.commit()
    return conteudo.id


def extrair_metadados_pdf(caminho):
    doc = fitz.open(caminho)
    meta = doc.metadata
    autor = meta.get('author')
    titulo = meta.get('title')
    # Tenta inferir tema pela primeira página
    tema = None
    if doc.page_count > 0:
        primeira_pagina = doc[0].get_text()
        linhas = primeira_pagina.split('\n')
        if linhas:
            tema = linhas[0] if len(linhas[0]) < 80 else None
    return autor, tema, titulo


def extrair_metadados_docx(caminho):
    doc = docx.Document(caminho)
    props = doc.core_properties
    autor = props.author
    titulo = props.title
    # Tenta inferir tema pela primeira linha
    tema = None
    if doc.paragraphs:
        primeira_linha = doc.paragraphs[0].text
        tema = primeira_linha if len(primeira_linha) < 80 else None
    return autor, tema, titulo


def extrair_metadados_txt(caminho):
    with open(caminho, "r", encoding="utf-8") as f:
        linhas = [linha.strip() for linha in f.readlines() if linha.strip()]
    autor = None
    tema = None
    for linha in linhas[:5]:
        if linha.lower().startswith("autor:"):
            autor = linha.split(":", 1)[-1].strip()
        if not tema:
            tema = linha if len(linha) < 80 else None
    return autor, tema, None


def extrair_metadados_html(caminho):
    with open(caminho, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
        texto = soup.get_text().splitlines()
        linhas = [linha.strip() for linha in texto if linha.strip()]
    autor = None
    tema = None
    for linha in linhas[:5]:
        if linha.lower().startswith("autor:"):
            autor = linha.split(":", 1)[-1].strip()
        if not tema:
            tema = linha if len(linha) < 80 else None
    return autor, tema, None


def extrair_metadados_ppt(caminho):
    prs = Presentation(caminho)
    autor = None
    tema = None
    # Tenta pegar título do primeiro slide
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text") and shape.text:
                if not tema:
                    tema = shape.text if len(shape.text) < 80 else None
                if shape.text.lower().startswith("autor:"):
                    autor = shape.text.split(":", 1)[-1].strip()
    return autor, tema, None


def processar_arquivo(caminho, tipo=None, autor=None, tema=None, fonte=None):
    # Detecta tipo por extensão, se não informado
    if not tipo:
        _, ext = os.path.splitext(caminho)
        tipo = ext.lower().lstrip(".")

    if tipo == "pdf":
        texto = extrair_pdf(caminho)
        auto_autor, auto_tema, auto_titulo = extrair_metadados_pdf(caminho)
    elif tipo == "docx" or tipo == "doc":
        texto = extrair_docx(caminho)
        auto_autor, auto_tema, auto_titulo = extrair_metadados_docx(caminho)
    elif tipo == "html":
        texto = extrair_html(caminho)
        auto_autor, auto_tema, auto_titulo = extrair_metadados_html(caminho)
    elif tipo == "txt":
        texto = extrair_txt(caminho)
        auto_autor, auto_tema, auto_titulo = extrair_metadados_txt(caminho)
    elif tipo == "ppt":
        texto = extrair_ppt(caminho)
        auto_autor, auto_tema, auto_titulo = extrair_metadados_ppt(caminho)
    else:
        raise ValueError("Tipo de arquivo não suportado")
    titulo = auto_titulo or os.path.basename(caminho)
    autor = autor or auto_autor
    tema = tema or auto_tema
    try:
        return salvar_conteudo(titulo, texto, tipo, autor, tema, fonte)
    except OperationalError:
        # Se o DB cair no meio, tenta fallback para SQLite
        try:
            global engine, session
            engine = create_engine("sqlite:////data/app.db")
            SessionLocal = sessionmaker(bind=engine)
            session = SessionLocal()
            Base.metadata.create_all(engine)
            return salvar_conteudo(titulo, texto, tipo, autor, tema, fonte)
        except Exception as e:
            raise e


def indexar_conteudo_teologico():
    """
    Indexa todos os textos do banco de dados no ChromaDB usando embeddings
    Ollama.
    """
    embedding_model = OllamaEmbeddings(model="gemma:2b")
    db = Chroma(
        persist_directory="./chroma_db",
        embedding_function=embedding_model
    )
    conteudos = session.query(ConteudoTeologico).all()
    documentos = []
    metadados = []
    for c in conteudos:
        documentos.append(c.texto)
        metadados.append({
            "id": c.id,
            "titulo": c.titulo,
            "tipo": c.tipo,
            "autor": c.autor,
            "tema": c.tema,
            "fonte": c.fonte
        })
    if documentos:
        db.add_texts(documentos, metadados)
        print(f"Indexados {len(documentos)} documentos no ChromaDB.")
    else:
        print("Nenhum conteúdo encontrado para indexar.")


if __name__ == "__main__":
    # Exemplo: indexar todo conteúdo teológico ao iniciar o script
    indexar_conteudo_teologico()

# Exemplo de uso:
# processar_arquivo(
#     "/caminho/para/arquivo.pdf",
#     "pdf",
#     autor="Desconhecido",
#     tema="Graça",
#     fonte="Acervo pessoal"
# )
# indexar_conteudo_teologico()
