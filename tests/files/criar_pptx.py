from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Sola Scriptura"
content = slide.shapes.add_textbox(100, 150, 500, 300)
content.text = "A autoridade final da fé cristã está nas Escrituras."
prs.save('j:/Projetos Python/eklesia.ia.python/tests/files/teste.pptx')
